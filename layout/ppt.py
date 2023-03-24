from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_ANCHOR,PP_ALIGN
from pptx.text.text import TextFrame
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
import lorem
import random


class PPTPage:
    def __init__(self, width=Inches(8.27), height=Inches(11.69)):
        # 创建一个空白的ppt文件
        self.prs = Presentation()
        # 设置幻灯片的尺寸，默认为A4纸的尺寸
        self.prs.slide_width = width
        self.prs.slide_height = height
        # 添加一个空白的幻灯片
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def add_text(self, x, y, w, h):
        # 输入相对位置，x,y,w,h，单位为百分比，范围为0到1
        # 计算绝对位置，单位为英寸，范围为0到幻灯片宽度或高度
        left = x * self.prs.slide_width
        top = y * self.prs.slide_height
        width = w * self.prs.slide_width
        height = h * self.prs.slide_height

        # 在幻灯片中添加一个文本框，并设置自动调整大小和正常的文本对齐方式
        textbox = self.slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 设置文本框的word_wrap属性为True以实现自动换行
        textbox.text_frame.word_wrap = True

        # 使用Lorem库生成占位文本并将其添加到文本框中
        
        p=textbox.text_frame.paragraphs[0]
        
        # 设置段落字体大小 
        font_size=Pt(8)
        p.font.size=font_size
        
        # 根据文本框的大小和字体大小，估计其可容纳的字符数
        char_per_line = max(1,int(width / (font_size/(1.4*2))))
        lines_per_box = max(1,int(height / (font_size/0.85)))
        max_chars = char_per_line * lines_per_box

        s="      "if lines_per_box>1 else ""
        while True:
        # 以单词为单位进行填充
            sentence=lorem.sentence()
            words = sentence.split()
            for word in words:
                word = " ".join(word)
                if len(s + word) < max_chars:
                    s += word + '   '
                else:
                    break
            else:
                continue
            break
        s+="."
        p.text=s
    
    def add_title(self, x, y, w, h):
        # 输入相对位置，x,y,w,h，单位为百分比，范围为0到1
        # 计算绝对位置，单位为英寸，范围为0到幻灯片宽度或高度
        left = x * self.prs.slide_width
        top = y * self.prs.slide_height
        width = w * self.prs.slide_width
        height = h * self.prs.slide_height

        # 在幻灯片中添加一个文本框，并设置自动调整大小和正常的文本对齐方式
        textbox = self.slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 设置文本框的word_wrap属性为True以实现自动换行
        #textbox.text_frame.word_wrap = True

        # 使用Lorem库生成占位文本并将其添加到文本框中
        
        p=textbox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        # 设置段落字体大小 
        font_size=max(Pt(10),height*0.7)
        p.font.size=font_size
        p.font.bold = True
        
        # 根据文本框的大小和字体大小，估计其可容纳的字符数
        char_per_line = max(1,int(width / (font_size/(1.6*2))))
        lines_per_box = 1
        max_chars = char_per_line * lines_per_box

        s="      "if lines_per_box>1 else ""
        while True:
        # 以单词为单位进行填充
            sentence=lorem.sentence()
            words = sentence.split()
            for word in words:
                word = " ".join(word)
                if len(s + word) < max_chars:
                    s += word + '   '
                else:
                    break
            else:
                continue
            break
        p.text=s
            
    def add_list(self, x, y, w, h):
        # 输入相对位置，x,y,w,h，单位为百分比，范围为0到1
        # 计算绝对位置，单位为英寸，范围为0到幻灯片宽度或高度
        left = x * self.prs.slide_width
        top = y * self.prs.slide_height
        width = w * self.prs.slide_width
        height = h * self.prs.slide_height

        # 在幻灯片中添加一个文本框，并设置自动调整大小和正常的文本对齐方式
        textbox = self.slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 设置文本框的word_wrap属性为True以实现自动换行
        textbox.text_frame.word_wrap = True

        # 设置段落字体大小 
        font_size=Pt(8)

        # 根据文本框的大小和字体大小，估计其可容纳的字符数
        char_per_line = max(1,int(width / (font_size/(1.45*2.2))))
        lines_per_box = max(1,int(height / (font_size/0.88)))

        lines=0

        sentence=" ".join(lorem.sentence())
        p=textbox.text_frame.paragraphs[0]
        p.font.size=font_size
        p.text=sentence
        lines+=(len(sentence)+char_per_line-1)/char_per_line
        index=0
        while True:
            if lines>=lines_per_box:break
            sentence=" ".join(lorem.sentence())
            p=textbox.text_frame.add_paragraph()
            p.font.size=font_size
            p.font.bold = True
            index+=1
            p.text=f"{index}."
            lines+=1

            p=textbox.text_frame.add_paragraph()
            p.font.size=font_size
            p.text=sentence
            p.level=1
            lines+=(len(sentence)+char_per_line-1)/char_per_line
            
    def add_table(self, x, y, w, h):
        # 输入相对位置，x,y,w,h，单位为百分比，范围为0到1
        # 计算绝对位置，单位为英寸，范围为0到幻灯片宽度或高度
        left =int( x * self.prs.slide_width)
        top = int(y * self.prs.slide_height)
        width = int(w * self.prs.slide_width)
        height = int(h * self.prs.slide_height)
        #self.slide.shapes.add_picture("save/table.png", left, top, width=width,height=height)
        cols = max(1,int(width / (Pt(8)*7)))
        rows = max(1,int(height / (Pt(8)*2.5)))

        table = self.slide.shapes.add_table(rows, cols, left, top, width, height).table
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                cell.fill.solid()
                if i==0:
                    cell.fill.fore_color.rgb=RGBColor(200, 200, 200)
                else:
                    cell.fill.fore_color.rgb=RGBColor(250, 250, 250)
                # 插入文字
                text_frame = cell.text_frame
                text_frame.clear() # 清除默认的段落
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE # 设置垂直居中
                p = text_frame.paragraphs[0] # 创建新的段落
                # 设置文本内容
                if j==0 and i==0:
                    p.text=""
                elif i==0:
                    p.text=f"table row {j}"
                elif j==0:
                    p.text=f"col {i}"
                else:
                    p.text=f"{random.random():.5f}"
                p.alignment = PP_ALIGN.CENTER # 设置水平居中
                p.font.size = Pt(8) # 设置字体大小为10磅
                p.font.color.rgb = RGBColor(0, 0, 0) # 设置字体颜色为黑色
                


    def add_figure(self, x, y, w, h):
        # 输入相对位置，x,y,w,h，单位为百分比，范围为0到1
        # 计算绝对位置，单位为英寸，范围为0到幻灯片宽度或高度
        left =int( x * self.prs.slide_width)
        top = int(y * self.prs.slide_height)
        width = int(w * self.prs.slide_width)
        height = int(h * self.prs.slide_height)
        self.slide.shapes.add_picture("../save/figure.png", left, top, width=width,height=height)

    def rander(self,data,size=2**8):
        m=size
        for item in data:
            if item["C"]=="text":
                self.add_text(item["X"]/m,item["Y"]/m,item["W"]/m,item["H"]/m)
            elif item["C"]=="title":
                self.add_title(item["X"]/m,item["Y"]/m,item["W"]/m,item["H"]/m)
            elif item["C"]=="list":
                self.add_list(item["X"]/m,item["Y"]/m,item["W"]/m,item["H"]/m)
            elif item["C"]=="table":
                self.add_table(item["X"]/m,item["Y"]/m,item["W"]/m,item["H"]/m)
            elif item["C"]=="figure":
                self.add_figure(item["X"]/m,item["Y"]/m,item["W"]/m,item["H"]/m)
        

    
    def save(self, filename="ppt.pptx"):
      # 保存这页ppt以单独的ppt文件格式保存，输入文件名（包括后缀名）
      self.prs.save(filename) 
    

if __name__=="__main__":
    page=PPTPage()
    m=2**8
    data=[{"C": "text", "X": 23, "Y": 27, "W": 93, "H": 3, "A": 17, "R": 7}, {"C": "table", "X": 23, "Y": 33, "W": 209, "H": 105, "A": 148, "R": 94}, {"C": "text", "X": 25, "Y": 137, "W": 101, "H": 3, "A": 17, "R": 6}, {"C": "text", "X": 23, "Y": 145, "W": 101, "H": 11, "A": 34, "R": 23}, {"C": "text", "X": 23, "Y": 160, "W": 101, "H": 30, "A": 55, "R": 60}, {"C": "title", "X": 23, "Y": 194, "W": 22, "H": 4, "A": 9, "R": 39}, {"C": "text", "X": 23, "Y": 198, "W": 101, "H": 19, "A": 44, "R": 39}, {"C": "text", "X": 23, "Y": 221, "W": 101, "H": 15, "A": 39, "R": 31}, {"C": "text", "X": 131, "Y": 145, "W": 101, "H": 68, "A": 83, "R": 117}, {"C": "text", "X": 131, "Y": 217, "W": 101, "H": 19, "A": 44, "R": 39}]
    page.rander(data,vocab_size=m)
    page.save()